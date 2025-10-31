# -*- coding: utf-8 -*-
"""
PowerPoint Document Processor

Extracts text content from PowerPoint presentations (.pptx and .ppt files)
and converts them to LangChain Document objects for vectorstore.
"""

import os
from typing import List, Optional
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text content from a PowerPoint file.
    
    Args:
        pptx_path: Path to the .pptx file
        
    Returns:
        Extracted text content
    """
    try:
        from pptx import Presentation
        
        # Load presentation
        prs = Presentation(pptx_path)
        
        text_content = []
        
        # Extract text from each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"\n--- Slide {slide_num} ---\n")
            
            # Extract text from shapes (text boxes, titles, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            # Extract notes (if any)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                    slide_text.append(f"\n[Notes: {notes_slide.notes_text_frame.text}]")
            
            # Add slide content to overall text
            if slide_text:
                text_content.extend(slide_text)
        
        # Join all text
        full_text = "\n".join(text_content)
        
        return full_text
    
    except ImportError:
        print("[ERROR] python-pptx library not installed. Run: pip install python-pptx")
        return ""
    except Exception as e:
        print(f"[ERROR] Failed to extract text from {pptx_path}: {e}")
        return ""


def process_powerpoint_file(
    pptx_path: str, 
    metadata: Optional[dict] = None
) -> List[Document]:
    """
    Process a PowerPoint file and return LangChain Documents.
    
    Args:
        pptx_path: Path to the PowerPoint file
        metadata: Additional metadata to attach to documents
        
    Returns:
        List of Document objects with chunked content
    """
    print(f"[*] Processing PowerPoint: {os.path.basename(pptx_path)}")
    
    # Extract text from PowerPoint
    text = extract_text_from_pptx(pptx_path)
    
    if not text or not text.strip():
        print(f"[WARNING] No text extracted from {pptx_path}")
        return []
    
    print(f"[OK] Extracted {len(text)} characters from PowerPoint")
    
    # Initialize text splitter for chunking
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=300,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    
    # Split text into chunks
    chunks = text_splitter.split_text(text)
    
    # Create Document objects
    documents = []
    for i, chunk in enumerate(chunks):
        # Create base metadata
        doc_metadata = {
            "source_type": "sharepoint_document",
            "source": "cloudfuze_sharepoint",
            "file_type": "pptx",
            "chunk_index": i,
            "total_chunks": len(chunks),
        }
        
        # Add file-specific metadata if provided
        if metadata:
            doc_metadata.update(metadata)
        
        # Create document
        doc = Document(
            page_content=chunk,
            metadata=doc_metadata
        )
        documents.append(doc)
    
    print(f"[OK] Created {len(documents)} document chunks from PowerPoint")
    
    return documents


def process_powerpoint_directory(pptx_directory: str) -> List[Document]:
    """
    Process all PowerPoint files in a directory.
    
    Args:
        pptx_directory: Path to directory containing PowerPoint files
        
    Returns:
        List of Document objects from all PowerPoint files
    """
    if not os.path.exists(pptx_directory):
        print(f"[ERROR] Directory not found: {pptx_directory}")
        return []
    
    print(f"[*] Processing PowerPoint files from: {pptx_directory}")
    
    documents = []
    pptx_files = []
    
    # Find all PowerPoint files
    for root, dirs, files in os.walk(pptx_directory):
        for file in files:
            if file.lower().endswith(('.pptx', '.ppt')):
                pptx_path = os.path.join(root, file)
                pptx_files.append(pptx_path)
    
    print(f"[*] Found {len(pptx_files)} PowerPoint files")
    
    # Process each PowerPoint file
    for pptx_path in pptx_files:
        try:
            # Create metadata with file information
            metadata = {
                "file_name": os.path.basename(pptx_path),
                "file_path": pptx_path,
            }
            
            # Process file
            file_docs = process_powerpoint_file(pptx_path, metadata)
            documents.extend(file_docs)
            
        except Exception as e:
            print(f"[ERROR] Failed to process {pptx_path}: {e}")
            continue
    
    print(f"[OK] Processed {len(pptx_files)} PowerPoint files")
    print(f"[OK] Total documents created: {len(documents)}")
    
    return documents


def chunk_powerpoint_documents(
    documents: List[Document], 
    chunk_size: int = 1500, 
    chunk_overlap: int = 300
) -> List[Document]:
    """
    Re-chunk PowerPoint documents with custom parameters.
    
    Args:
        documents: List of Document objects
        chunk_size: Maximum chunk size in characters
        chunk_overlap: Overlap between chunks
        
    Returns:
        List of re-chunked Document objects
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    
    # Split documents
    chunked_docs = text_splitter.split_documents(documents)
    
    # Update chunk indices
    for i, doc in enumerate(chunked_docs):
        doc.metadata['chunk_index'] = i
        doc.metadata['total_chunks'] = len(chunked_docs)
    
    return chunked_docs

